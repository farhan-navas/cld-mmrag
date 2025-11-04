import logging
from typing import List, Dict, Any
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger("indexer.extractors.pptx")

def extract_slides(path: Path) -> List[Dict[str, Any]]:
    """
    extract slides from a pptx file
    """
    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.warning(f"Failed to load PowerPoint file {path.name}: {e}")
        return []
    
    slides_data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_content = {
            "slide_index": i,
            "title": None,
            "text_boxes": [],
            "tables": [],
            "notes": None
        }
        
        # extract title
        if slide.shapes.title and slide.shapes.title.text:
            slide_content["title"] = slide.shapes.title.text.strip()
        
        # extract text from all shapes
        for shp in slide.shapes:
            if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                text = shp.text_frame.text or ""  # type: ignore[attr-defined]
                text = text.strip()
                if text and text != slide_content["title"]:
                    slide_content["text_boxes"].append(text)
            
            # Tables
            if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    table = shp.table  # type: ignore[attr-defined]
                    rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows.append(row_data)
                    if rows:
                        slide_content["tables"].append(rows)
                except Exception as e:
                    logger.warning(f"Failed to extract table from slide {i}: {e}")
        
        # extract speaker notes for additional context
        if slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text or ""
                    notes = notes.strip()
                    if notes:
                        slide_content["notes"] = notes
            except Exception as e:
                logger.warning(f"Failed to extract notes from slide {i}: {e}")
        
        slides_data.append(slide_content)
    
    return slides_data
