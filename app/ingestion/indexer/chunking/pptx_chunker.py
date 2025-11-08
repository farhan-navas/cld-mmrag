"""
PowerPoint chunking strategy.

This module handles chunking of PowerPoint files slide-by-slide,
including titles, text boxes, speaker notes, and embedded tables.
"""

import logging
import json
from typing import List, Dict, Any
from pathlib import Path

from app.config import config
from app.ingestion.indexer.chunking.text_splitter import recursive_split_text

logger = logging.getLogger("indexer.chunking.pptx")


def chunk_pptx(path: Path, title: str, filepath: str) -> List[Dict[str, Any]]:
    """
    Chunk PowerPoint files slide-by-slide, including title, text boxes, 
    speaker notes, and embedded tables.
    
    :param path: Path to the PowerPoint file
    :param title: Document title
    :param filepath: Filepath for metadata
    :return: List of chunk dictionaries
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.error("python-pptx is required for PowerPoint processing. Install with: pip install python-pptx")
        return []
    
    chunks = []
    max_chars = config.indexing.chunk_max_chars
    include_notes = config.indexing.pptx_include_notes
    
    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.warning(f"Failed to load PowerPoint file {path.name}: {e}")
        return []
    
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        
        # Extract slide title
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(f"# {slide.shapes.title.text.strip()}")
        
        # Extract text from all shapes
        for shp in slide.shapes:
            # Text frames
            if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                text = shp.text_frame.text or "" # pyright: ignore[reportAttributeAccessIssue]
                text = text.strip()
                if text and text not in parts:  # avoid duplicating title
                    parts.append(text)
            
            # Tables
            if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    table = shp.table # pyright: ignore[reportAttributeAccessIssue]
                    rows = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(row_data) + " |")
                    
                    if rows:
                        # Add markdown table header separator after first row
                        if len(rows) > 1:
                            rows.insert(1, "| " + " | ".join(["---"] * len(table.columns)) + " |")
                        table_md = "\n".join(rows)
                        parts.append(f"\n**Table:**\n{table_md}\n")
                except Exception as e:
                    logger.warning(f"Failed to extract table from slide {i}: {e}")
        
        # Extract speaker notes if enabled
        if include_notes and slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes = notes_slide.notes_text_frame.text or ""
                    notes = notes.strip()
                    if notes:
                        parts.append(f"\n**Speaker Notes:**\n{notes}")
            except Exception as e:
                logger.warning(f"Failed to extract notes from slide {i}: {e}")
        
        # Combine all parts
        body = "\n\n".join(parts).strip()
        if not body:
            continue
        
        # Enforce chunk size
        if len(body) > max_chars:
            for tc in recursive_split_text(body, max_chars):
                chunks.append({
                    "title": title,
                    "filepath": filepath,
                    "page": i,  # use slide index as page number
                    "section_path": f"Slide {i}",
                    "content": tc,
                    "content_markdown": tc,
                    "bbox": None,
                    "metadata_json": json.dumps({
                        "filetype": "pptx",
                        "slide_index": i
                    })
                })
        else:
            chunks.append({
                "title": title,
                "filepath": filepath,
                "page": i,
                "section_path": f"Slide {i}",
                "content": body,
                "content_markdown": body,
                "bbox": None,
                "metadata_json": json.dumps({
                    "filetype": "pptx",
                    "slide_index": i
                })
            })
    
    return chunks
